import os
from typing import Dict, Any
from pptx import Presentation
import pymupdf


def parse_pptx(file_path: str) -> Dict[str, Any]:
    """Extract text from PowerPoint files."""
    prs = Presentation(file_path)
    slides_text = []
    
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = f"Slide {i}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = getattr(shape, "text", "")
                if text:
                    slide_text += text + "\n"
        slides_text.append(slide_text)
    
    return {
        "content": "\n\n".join(slides_text),
        "slide_count": len(prs.slides)
    }


def parse_pdf(file_path: str) -> Dict[str, Any]:
    """Extract text from PDF files."""
    doc = pymupdf.open(file_path)
    pages_text = []
    page_count = len(doc)
    
    for i in range(page_count):
        page = doc[i]
        text = page.get_text()  # type: ignore
        if text.strip():
            pages_text.append(f"Page {i+1}:\n{text}")
    
    doc.close()
    
    return {
        "content": "\n\n".join(pages_text),
        "page_count": page_count
    }


def parse_document(file_path: str) -> Dict[str, Any]:
    """Parse document based on file extension."""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if file_ext == '.pptx':
        return parse_pptx(file_path)
    elif file_ext == '.pdf':
        return parse_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_ext}. Only PDF and PPTX files are supported.")
